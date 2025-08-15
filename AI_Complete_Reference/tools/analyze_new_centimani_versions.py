#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
新版本Centimani文件分析工具

用於分析2023年平台培訓的新版本文件，包括：
1. 新版本PPT (20210217)
2. 新版本儀器參數 (Ver1.71)
3. 新版本腳本用戶指南
4. iPLAS參數列表

作者: AI Assistant
版本: 1.0.0
更新日期: 2025-08-11
"""

import os
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

class NewVersionAnalyzer:
    """新版本Centimani文件分析器"""
    
    def __init__(self):
        self.analysis_result = {
            "analysis_time": datetime.now().isoformat(),
            "new_versions": {},
            "version_comparison": {},
            "new_features": [],
            "updated_rules": [],
            "recommendations": []
        }
    
    def analyze_new_versions(self, base_path: str):
        """分析新版本文件"""
        print("🔍 開始分析新版本Centimani文件...")
        
        # 分析新版本PPT
        ppt_path = os.path.join(base_path, "Centimani - How to edit scripts (20210217).pptx")
        if os.path.exists(ppt_path):
            self._analyze_new_ppt(ppt_path)
        
        # 分析新版本儀器參數
        instrument_path = os.path.join(base_path, "Centimani instrument parameter -Ver1.71(20230411).xlsx")
        if os.path.exists(instrument_path):
            self._analyze_new_instrument_params(instrument_path)
        
        # 分析新版本腳本用戶指南
        script_guide_path = os.path.join(base_path, "Centimani Script User Guide(202230104).xlsx")
        if os.path.exists(script_guide_path):
            self._analyze_new_script_guide(script_guide_path)
        
        # 分析iPLAS參數列表
        iplas_path = os.path.join(base_path, "dbaccess2_iPLAS_Argument_List_v2.0.xlsx")
        if os.path.exists(iplas_path):
            self._analyze_iplas_params(iplas_path)
        
        # 生成版本比較
        self._generate_version_comparison()
        
        # 生成更新建議
        self._generate_update_recommendations()
        
        print("✅ 新版本分析完成！")
        return self.analysis_result
    
    def _analyze_new_ppt(self, ppt_path: str):
        """分析新版本PPT"""
        if not PPTX_AVAILABLE:
            print("⚠️ 無法分析PPT文件，python-pptx未安裝")
            return
        
        try:
            print(f"📊 分析新版本PPT: {os.path.basename(ppt_path)}")
            prs = Presentation(ppt_path)
            
            self.analysis_result["new_versions"]["ppt"] = {
                "file_name": os.path.basename(ppt_path),
                "file_size": os.path.getsize(ppt_path),
                "total_slides": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
                "slides_content": []
            }
            
            # 分析每個幻燈片
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = self._extract_slide_text(slide)
                if slide_text:
                    self.analysis_result["new_versions"]["ppt"]["slides_content"].append({
                        "slide": slide_num,
                        "text": slide_text
                    })
            
            print(f"   - 總幻燈片數: {len(prs.slides)}")
            print(f"   - 文字內容數: {len(self.analysis_result['new_versions']['ppt']['slides_content'])}")
            
        except Exception as e:
            print(f"❌ PPT分析失敗: {e}")
    
    def _analyze_new_instrument_params(self, excel_path: str):
        """分析新版本儀器參數"""
        if not PANDAS_AVAILABLE:
            print("⚠️ 無法分析Excel文件，pandas未安裝")
            return
        
        try:
            print(f"📊 分析新版本儀器參數: {os.path.basename(excel_path)}")
            
            # 嘗試讀取Excel文件
            excel_file = pd.ExcelFile(excel_path)
            sheet_names = excel_file.sheet_names
            
            self.analysis_result["new_versions"]["instrument_params"] = {
                "file_name": os.path.basename(excel_path),
                "file_size": os.path.getsize(excel_path),
                "sheet_names": sheet_names,
                "sheets_info": {}
            }
            
            # 分析每個工作表
            for sheet_name in sheet_names:
                try:
                    df = pd.read_excel(excel_path, sheet_name=sheet_name)
                    self.analysis_result["new_versions"]["instrument_params"]["sheets_info"][sheet_name] = {
                        "rows": len(df),
                        "columns": len(df.columns),
                        "column_names": df.columns.tolist(),
                        "sample_data": df.head(3).to_dict('records') if len(df) > 0 else []
                    }
                except Exception as e:
                    print(f"   - 工作表 {sheet_name} 分析失敗: {e}")
            
            print(f"   - 工作表數: {len(sheet_names)}")
            print(f"   - 工作表: {', '.join(sheet_names)}")
            
        except Exception as e:
            print(f"❌ 儀器參數分析失敗: {e}")
    
    def _analyze_new_script_guide(self, excel_path: str):
        """分析新版本腳本用戶指南"""
        if not PANDAS_AVAILABLE:
            print("⚠️ 無法分析Excel文件，pandas未安裝")
            return
        
        try:
            print(f"📊 分析新版本腳本用戶指南: {os.path.basename(excel_path)}")
            
            excel_file = pd.ExcelFile(excel_path)
            sheet_names = excel_file.sheet_names
            
            self.analysis_result["new_versions"]["script_guide"] = {
                "file_name": os.path.basename(excel_path),
                "file_size": os.path.getsize(excel_path),
                "sheet_names": sheet_names,
                "sheets_info": {}
            }
            
            # 分析每個工作表
            for sheet_name in sheet_names:
                try:
                    df = pd.read_excel(excel_path, sheet_name=sheet_name)
                    self.analysis_result["new_versions"]["script_guide"]["sheets_info"][sheet_name] = {
                        "rows": len(df),
                        "columns": len(df.columns),
                        "column_names": df.columns.tolist(),
                        "sample_data": df.head(3).to_dict('records') if len(df) > 0 else []
                    }
                except Exception as e:
                    print(f"   - 工作表 {sheet_name} 分析失敗: {e}")
            
            print(f"   - 工作表數: {len(sheet_names)}")
            print(f"   - 工作表: {', '.join(sheet_names)}")
            
        except Exception as e:
            print(f"❌ 腳本用戶指南分析失敗: {e}")
    
    def _analyze_iplas_params(self, excel_path: str):
        """分析iPLAS參數列表"""
        if not PANDAS_AVAILABLE:
            print("⚠️ 無法分析Excel文件，pandas未安裝")
            return
        
        try:
            print(f"📊 分析iPLAS參數列表: {os.path.basename(excel_path)}")
            
            excel_file = pd.ExcelFile(excel_path)
            sheet_names = excel_file.sheet_names
            
            self.analysis_result["new_versions"]["iplas_params"] = {
                "file_name": os.path.basename(excel_path),
                "file_size": os.path.getsize(excel_path),
                "sheet_names": sheet_names,
                "sheets_info": {}
            }
            
            # 分析每個工作表
            for sheet_name in sheet_names:
                try:
                    df = pd.read_excel(excel_path, sheet_name=sheet_name)
                    self.analysis_result["new_versions"]["iplas_params"]["sheets_info"][sheet_name] = {
                        "rows": len(df),
                        "columns": len(df.columns),
                        "column_names": df.columns.tolist(),
                        "sample_data": df.head(3).to_dict('records') if len(df) > 0 else []
                    }
                except Exception as e:
                    print(f"   - 工作表 {sheet_name} 分析失敗: {e}")
            
            print(f"   - 工作表數: {len(sheet_names)}")
            print(f"   - 工作表: {', '.join(sheet_names)}")
            
        except Exception as e:
            print(f"❌ iPLAS參數分析失敗: {e}")
    
    def _extract_slide_text(self, slide) -> str:
        """提取幻燈片文字內容"""
        text_content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
        
        return " | ".join(text_content) if text_content else ""
    
    def _generate_version_comparison(self):
        """生成版本比較"""
        print("🔄 生成版本比較...")
        
        # 比較PPT版本
        if "ppt" in self.analysis_result["new_versions"]:
            old_ppt_info = {
                "version": "v1.4(20180427)",
                "total_slides": 19
            }
            
            new_ppt_info = self.analysis_result["new_versions"]["ppt"]
            
            self.analysis_result["version_comparison"]["ppt"] = {
                "old_version": old_ppt_info,
                "new_version": {
                    "version": "(20210217)",
                    "total_slides": new_ppt_info.get("total_slides", 0)
                },
                "changes": {
                    "slide_count_change": new_ppt_info.get("total_slides", 0) - old_ppt_info["total_slides"],
                    "update_date": "2021-02-17 (vs 2018-04-27)"
                }
            }
        
        # 比較儀器參數版本
        if "instrument_params" in self.analysis_result["new_versions"]:
            old_instrument_info = {
                "version": "Ver1.65 (20230104)",
                "update_date": "2023-01-04"
            }
            
            new_instrument_info = {
                "version": "Ver1.71(20230411)",
                "update_date": "2023-04-11"
            }
            
            self.analysis_result["version_comparison"]["instrument_params"] = {
                "old_version": old_instrument_info,
                "new_version": new_instrument_info,
                "changes": {
                    "version_increment": "1.65 → 1.71",
                    "update_date": "2023-04-11 (vs 2023-01-04)",
                    "time_gap": "3個月"
                }
            }
    
    def _generate_update_recommendations(self):
        """生成更新建議"""
        print("💡 生成更新建議...")
        
        recommendations = []
        
        # 基於版本比較的建議
        if "ppt" in self.analysis_result["version_comparison"]:
            ppt_comp = self.analysis_result["version_comparison"]["ppt"]
            if ppt_comp["changes"]["slide_count_change"] != 0:
                recommendations.append({
                    "type": "ppt_update",
                    "priority": "high",
                    "description": f"PPT版本已更新至{ppt_comp['new_version']['version']}，建議更新速查表內容",
                    "action": "分析新版本PPT的差異，更新速查表"
                })
        
        if "instrument_params" in self.analysis_result["version_comparison"]:
            inst_comp = self.analysis_result["version_comparison"]["instrument_params"]
            recommendations.append({
                "type": "instrument_update",
                "priority": "medium",
                "description": f"儀器參數已更新至{inst_comp['new_version']['version']}，包含新功能和參數",
                "action": "檢查新版本儀器參數，更新相關設定說明"
            })
        
        # 基於新功能的建議
        if "iplas_params" in self.analysis_result["new_versions"]:
            recommendations.append({
                "type": "new_feature",
                "priority": "high",
                "description": "新增iPLAS參數列表v2.0，包含新的參數和功能",
                "action": "將iPLAS參數整合到速查表中，提供完整的參數參考"
            })
        
        # 通用建議
        recommendations.extend([
            {
                "type": "general",
                "priority": "medium",
                "description": "新版本文件包含更多實用資訊和最佳實踐",
                "action": "定期檢查和更新速查表，保持與最新版本同步"
            },
            {
                "type": "integration",
                "priority": "low",
                "description": "考慮將新版本內容與現有速查表整合",
                "action": "創建版本對照表，幫助用戶了解版本差異"
            }
        ])
        
        self.analysis_result["recommendations"] = recommendations
    
    def save_analysis_report(self, output_file: str = None) -> str:
        """保存分析報告"""
        if output_file is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = f"New_Version_Analysis_Report_{timestamp}.json"
        
        try:
            # 創建可序列化的副本
            serializable_result = self._make_serializable(self.analysis_result)
            
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(serializable_result, f, indent=2, ensure_ascii=False)
            
            print(f"📄 分析報告已保存: {output_file}")
            return output_file
            
        except Exception as e:
            print(f"❌ 保存分析報告失敗: {e}")
            return ""
    
    def _make_serializable(self, obj):
        """將對象轉換為可JSON序列化的格式"""
        if isinstance(obj, dict):
            return {k: self._make_serializable(v) for k, v in obj.items()}
        elif isinstance(obj, list):
            return [self._make_serializable(item) for item in obj]
        elif isinstance(obj, datetime):
            return obj.isoformat()
        else:
            return obj
    
    def generate_update_summary(self) -> str:
        """生成更新摘要"""
        summary = []
        summary.append("# Centimani 新版本分析摘要")
        summary.append(f"*分析時間: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*")
        summary.append("")
        
        # 版本比較摘要
        summary.append("## 🔄 版本比較")
        summary.append("")
        
        if "ppt" in self.analysis_result["version_comparison"]:
            ppt_comp = self.analysis_result["version_comparison"]["ppt"]
            summary.append(f"### PPT版本更新")
            summary.append(f"- **舊版本**: {ppt_comp['old_version']['version']} ({ppt_comp['old_version']['total_slides']} 頁)")
            summary.append(f"- **新版本**: {ppt_comp['new_version']['version']} ({ppt_comp['new_version']['total_slides']} 頁)")
            summary.append(f"- **更新日期**: {ppt_comp['changes']['update_date']}")
            summary.append("")
        
        if "instrument_params" in self.analysis_result["version_comparison"]:
            inst_comp = self.analysis_result["version_comparison"]["instrument_params"]
            summary.append(f"### 儀器參數更新")
            summary.append(f"- **舊版本**: {inst_comp['old_version']['version']}")
            summary.append(f"- **新版本**: {inst_comp['new_version']['version']}")
            summary.append(f"- **版本增量**: {inst_comp['changes']['version_increment']}")
            summary.append(f"- **更新間隔**: {inst_comp['changes']['time_gap']}")
            summary.append("")
        
        # 新功能摘要
        summary.append("## 🆕 新功能與改進")
        summary.append("")
        
        if "iplas_params" in self.analysis_result["new_versions"]:
            summary.append("### iPLAS參數列表 v2.0")
            summary.append("- 新增完整的iPLAS參數參考")
            summary.append("- 提供v2.0版本的參數說明")
            summary.append("- 包含更多實用參數和設定選項")
            summary.append("")
        
        # 更新建議
        summary.append("## 💡 更新建議")
        summary.append("")
        
        for rec in self.analysis_result["recommendations"]:
            priority_icon = "🔴" if rec["priority"] == "high" else "🟡" if rec["priority"] == "medium" else "🟢"
            summary.append(f"{priority_icon} **{rec['type'].upper()}**: {rec['description']}")
            summary.append(f"   - 建議行動: {rec['action']}")
            summary.append("")
        
        return "\n".join(summary)


def main():
    """主函數"""
    print("=" * 70)
    print("Centimani 新版本分析工具")
    print("=" * 70)
    
    # 檢查依賴套件
    if not PPTX_AVAILABLE:
        print("⚠️  警告: 未安裝python-pptx套件")
        print("   請執行: pip install python-pptx")
        print()
    
    if not PANDAS_AVAILABLE:
        print("⚠️  警告: 未安裝pandas套件")
        print("   請執行: pip install pandas")
        print()
    
    # 設定新版本文件路徑
    base_path = "../../Centimani DOC/CENTIMANIA/CENTIMANIA_2023_Platform_Training-20230718"
    
    if not os.path.exists(base_path):
        print(f"❌ 未找到新版本文件目錄: {base_path}")
        print("請確保文件位於正確位置")
        return
    
    try:
        # 創建分析器
        analyzer = NewVersionAnalyzer()
        
        # 分析新版本文件
        analysis_result = analyzer.analyze_new_versions(base_path)
        
        # 生成更新摘要
        print("\n生成更新摘要...")
        update_summary = analyzer.generate_update_summary()
        
        # 保存分析報告
        print("\n保存分析報告...")
        report_file = analyzer.save_analysis_report()
        
        # 顯示結果摘要
        print("\n" + "=" * 70)
        print("分析結果摘要")
        print("=" * 70)
        
        if "new_versions" in analysis_result:
            for version_type, info in analysis_result["new_versions"].items():
                if "file_name" in info:
                    print(f"✅ {version_type}: {info['file_name']}")
        
        print(f"📄 分析報告: {report_file}")
        print("\n✅ 新版本分析完成！")
        
        # 顯示更新摘要
        print("\n" + "=" * 70)
        print("更新摘要")
        print("=" * 70)
        print(update_summary)
        
    except Exception as e:
        print(f"❌ 分析失敗: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main() 