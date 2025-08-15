#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT分析工具

用於分析Centimani腳本編輯教學PPT，提取關鍵資訊並生成速查表
由於PPT包含大量圖片和描述，本工具將嘗試：
1. 提取文字內容
2. 分析幻燈片結構
3. 生成腳本編輯速查表
4. 記錄重要的教學要點

作者: AI Assistant
版本: 1.0.0
更新日期: 2024-12-19
"""

import os
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

class PPTAnalyzer:
    """PPT檔案分析器"""
    
    def __init__(self):
        self.analysis_result = {
            "file_info": {},
            "slides": [],
            "text_content": [],
            "key_points": [],
            "script_editing_guide": {},
            "analysis_time": datetime.now().isoformat()
        }
    
    def analyze_ppt_file(self, ppt_file_path: str) -> Dict[str, Any]:
        """分析PPT檔案"""
        if not PPTX_AVAILABLE:
            return self._create_fallback_analysis(ppt_file_path)
        
        try:
            print(f"正在分析PPT檔案: {ppt_file_path}")
            
            # 載入PPT
            prs = Presentation(ppt_file_path)
            
            # 基本檔案資訊
            self.analysis_result["file_info"] = {
                "file_name": os.path.basename(ppt_file_path),
                "file_size": os.path.getsize(ppt_file_path),
                "total_slides": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height
            }
            
            # 分析每個幻燈片
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_info = self._analyze_slide(slide, slide_num)
                self.analysis_result["slides"].append(slide_info)
                
                # 提取文字內容
                slide_text = self._extract_slide_text(slide)
                if slide_text:
                    self.analysis_result["text_content"].append({
                        "slide": slide_num,
                        "text": slide_text
                    })
            
            # 分析關鍵要點
            self._analyze_key_points()
            
            # 生成腳本編輯指南
            self._generate_script_editing_guide()
            
            print(f"PPT分析完成，共分析 {len(prs.slides)} 個幻燈片")
            return self.analysis_result
            
        except Exception as e:
            print(f"PPT分析失敗: {e}")
            return self._create_fallback_analysis(ppt_file_path)
    
    def _analyze_slide(self, slide, slide_num: int) -> Dict[str, Any]:
        """分析單個幻燈片"""
        slide_info = {
            "slide_number": slide_num,
            "shapes_count": len(slide.shapes),
            "text_boxes": [],
            "images": [],
            "tables": [],
            "charts": []
        }
        
        for shape in slide.shapes:
            shape_info = self._analyze_shape(shape)
            if shape_info:
                slide_info[shape_info["type"]].append(shape_info)
        
        return slide_info
    
    def _analyze_shape(self, shape) -> Optional[Dict[str, Any]]:
        """分析形狀物件"""
        try:
            if hasattr(shape, 'text') and shape.text.strip():
                return {
                    "type": "text_boxes",
                    "text": shape.text.strip(),
                    "position": (shape.left, shape.top),
                    "size": (shape.width, shape.height)
                }
            elif hasattr(shape, 'image'):
                return {
                    "type": "images",
                    "image_type": "embedded_image",
                    "position": (shape.left, shape.top),
                    "size": (shape.width, shape.height)
                }
            elif hasattr(shape, 'table'):
                return {
                    "type": "tables",
                    "rows": len(shape.table.rows),
                    "columns": len(shape.table.columns),
                    "position": (shape.left, shape.top)
                }
            elif hasattr(shape, 'chart'):
                return {
                    "type": "charts",
                    "chart_type": str(shape.chart.chart_type),
                    "position": (shape.left, shape.top)
                }
        except Exception:
            pass
        
        return None
    
    def _extract_slide_text(self, slide) -> str:
        """提取幻燈片文字內容"""
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text_parts.append(shape.text.strip())
        
        return " | ".join(text_parts) if text_parts else ""
    
    def _analyze_key_points(self):
        """分析關鍵要點"""
        key_points = []
        
        # 從文字內容中提取關鍵資訊
        for content in self.analysis_result["text_content"]:
            text = content["text"].lower()
            
            # 腳本編輯相關關鍵字
            if any(keyword in text for keyword in ["script", "edit", "create", "modify"]):
                key_points.append({
                    "slide": content["slide"],
                    "category": "script_editing",
                    "content": content["text"],
                    "importance": "high"
                })
            
            # 測試流程相關關鍵字
            if any(keyword in text for keyword in ["test", "flow", "procedure", "step"]):
                key_points.append({
                    "slide": content["slide"],
                    "category": "test_flow",
                    "content": content["text"],
                    "importance": "high"
                })
            
            # 儀器配置相關關鍵字
            if any(keyword in text for keyword in ["instrument", "device", "config", "setting"]):
                key_points.append({
                    "slide": content["slide"],
                    "category": "instrument_config",
                    "content": content["text"],
                    "importance": "medium"
                })
            
            # 錯誤處理相關關鍵字
            if any(keyword in text for keyword in ["error", "fail", "exception", "handle"]):
                key_points.append({
                    "slide": content["slide"],
                    "category": "error_handling",
                    "content": content["text"],
                    "importance": "medium"
                })
        
        self.analysis_result["key_points"] = key_points
    
    def _generate_script_editing_guide(self):
        """生成腳本編輯指南"""
        guide = {
            "basic_concepts": [],
            "editing_steps": [],
            "best_practices": [],
            "common_issues": [],
            "troubleshooting": []
        }
        
        # 根據關鍵要點分類整理
        for point in self.analysis_result["key_points"]:
            if "script_editing" in point["category"]:
                guide["editing_steps"].append({
                    "slide": point["slide"],
                    "description": point["content"]
                })
            elif "test_flow" in point["category"]:
                guide["basic_concepts"].append({
                    "slide": point["slide"],
                    "description": point["content"]
                })
            elif "instrument_config" in point["category"]:
                guide["best_practices"].append({
                    "slide": point["slide"],
                    "description": point["content"]
                })
            elif "error_handling" in point["category"]:
                guide["troubleshooting"].append({
                    "slide": point["slide"],
                    "description": point["content"]
                })
        
        self.analysis_result["script_editing_guide"] = guide
    
    def _create_fallback_analysis(self, ppt_file_path: str) -> Dict[str, Any]:
        """創建備用分析結果（當無法讀取PPT時）"""
        print("無法讀取PPT檔案，創建備用分析結果")
        
        return {
            "file_info": {
                "file_name": os.path.basename(ppt_file_path),
                "file_size": os.path.getsize(ppt_file_path),
                "file_type": "PowerPoint Presentation (.pptx)",
                "note": "需要安裝python-pptx套件來讀取PPT內容"
            },
            "slides": [],
            "text_content": [],
            "key_points": [],
            "script_editing_guide": {
                "basic_concepts": [],
                "editing_steps": [],
                "best_practices": [],
                "common_issues": [],
                "troubleshooting": []
            },
            "analysis_time": datetime.now().isoformat(),
            "status": "fallback_analysis"
        }
    
    def generate_quick_reference(self, output_file: str = None) -> str:
        """生成速查表"""
        if output_file is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = f"Centimani_Script_Editing_Quick_Reference_{timestamp}.md"
        
        try:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(self._format_quick_reference())
            
            print(f"速查表已生成: {output_file}")
            return output_file
            
        except Exception as e:
            print(f"生成速查表失敗: {e}")
            return ""
    
    def _format_quick_reference(self) -> str:
        """格式化速查表內容"""
        content = []
        
        # 標題
        content.append("# Centimani腳本編輯速查表")
        content.append("")
        content.append(f"**生成時間**: {self.analysis_result['analysis_time']}")
        content.append(f"**來源檔案**: {self.analysis_result['file_info'].get('file_name', 'Unknown')}")
        content.append("")
        
        # 檔案資訊
        content.append("## 📋 檔案資訊")
        content.append("")
        for key, value in self.analysis_result["file_info"].items():
            if key != "file_name":
                content.append(f"- **{key}**: {value}")
        content.append("")
        
        # 腳本編輯指南
        guide = self.analysis_result["script_editing_guide"]
        
        if guide["basic_concepts"]:
            content.append("## 🎯 基本概念")
            content.append("")
            for concept in guide["basic_concepts"]:
                content.append(f"- **幻燈片 {concept['slide']}**: {concept['description']}")
            content.append("")
        
        if guide["editing_steps"]:
            content.append("## ✏️ 編輯步驟")
            content.append("")
            for step in guide["editing_steps"]:
                content.append(f"- **幻燈片 {step['slide']}**: {step['description']}")
            content.append("")
        
        if guide["best_practices"]:
            content.append("## ✅ 最佳實踐")
            content.append("")
            for practice in guide["best_practices"]:
                content.append(f"- **幻燈片 {practice['slide']}**: {practice['description']}")
            content.append("")
        
        if guide["troubleshooting"]:
            content.append("## 🔧 故障排除")
            content.append("")
            for issue in guide["troubleshooting"]:
                content.append(f"- **幻燈片 {issue['slide']}**: {issue['description']}")
            content.append("")
        
        # 關鍵要點摘要
        if self.analysis_result["key_points"]:
            content.append("## 🔑 關鍵要點摘要")
            content.append("")
            for point in self.analysis_result["key_points"][:10]:  # 只顯示前10個
                content.append(f"- **{point['category']}** (幻燈片 {point['slide']}): {point['content'][:100]}...")
            content.append("")
        
        # 注意事項
        content.append("## ⚠️ 注意事項")
        content.append("")
        content.append("1. **圖片內容**: 本速查表主要基於文字內容生成，重要圖片請參考原始PPT")
        content.append("2. **完整資訊**: 建議結合原始PPT檔案使用，獲取完整的視覺資訊")
        content.append("3. **版本差異**: 不同版本的Centimani可能有差異，請以實際環境為準")
        content.append("4. **持續更新**: 根據使用經驗持續完善和更新速查表")
        content.append("")
        
        # 聯絡資訊
        content.append("## 📞 支援與聯絡")
        content.append("")
        content.append("- **技術支援**: 參考原始PPT檔案和Centimani官方文檔")
        content.append("- **問題回報**: 記錄具體問題和錯誤訊息")
        content.append("- **改進建議**: 提供使用體驗和改進建議")
        content.append("")
        
        return "\n".join(content)
    
    def save_analysis_report(self, output_file: str = None) -> str:
        """保存分析報告"""
        if output_file is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = f"PPT_Analysis_Report_{timestamp}.json"
        
        try:
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(self.analysis_result, f, indent=2, ensure_ascii=False)
            
            print(f"分析報告已保存: {output_file}")
            return output_file
            
        except Exception as e:
            print(f"保存分析報告失敗: {e}")
            return ""


def main():
    """主函數"""
    print("=" * 60)
    print("Centimani PPT分析工具")
    print("=" * 60)
    
    # 檢查依賴套件
    if not PPTX_AVAILABLE:
        print("⚠️  警告: 未安裝python-pptx套件")
        print("   請執行: pip install python-pptx")
        print("   或使用備用分析模式")
        print()
    
    # 查找PPT檔案
    ppt_file = "Centimani - How to edit scripts v1.4(20180427).pptx"
    ppt_path = f"../../Centimani DOC/{ppt_file}"
    
    if not os.path.exists(ppt_path):
        print(f"❌ 未找到PPT檔案: {ppt_path}")
        print("請確保PPT檔案位於正確位置")
        return
    
    try:
        # 創建分析器
        analyzer = PPTAnalyzer()
        
        # 分析PPT檔案
        print(f"開始分析PPT檔案: {ppt_file}")
        analysis_result = analyzer.analyze_ppt_file(ppt_path)
        
        # 生成速查表
        print("\n生成腳本編輯速查表...")
        quick_ref_file = analyzer.generate_quick_reference()
        
        # 保存分析報告
        print("\n保存分析報告...")
        report_file = analyzer.save_analysis_report()
        
        # 顯示結果摘要
        print("\n" + "=" * 60)
        print("分析結果摘要")
        print("=" * 60)
        
        if "total_slides" in analysis_result["file_info"]:
            print(f"總幻燈片數: {analysis_result['file_info']['total_slides']}")
            print(f"文字內容數: {len(analysis_result['text_content'])}")
            print(f"關鍵要點數: {len(analysis_result['key_points'])}")
        
        print(f"速查表: {quick_ref_file}")
        print(f"分析報告: {report_file}")
        
        print("\n✅ 分析完成！")
        
    except Exception as e:
        print(f"❌ 分析失敗: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main() 